import sys
import os
import time
import re
import zipfile
from datetime import datetime
import numpy as np

# GUI Imports
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QVBoxLayout, QHBoxLayout, 
    QWidget, QPushButton, QLabel, QFileDialog, QMessageBox,
    QFrame, QFormLayout, QDoubleSpinBox, QDateTimeEdit,
    QGraphicsDropShadowEffect, QSizePolicy, QDialog, QLineEdit, QTextEdit,
    QTextBrowser, QGraphicsProxyWidget, QMenu, QSlider, QCheckBox
)
from PyQt5.QtCore import QThread, pyqtSignal, Qt, QRectF, QPointF, QTimer, QPropertyAnimation
from PyQt5.QtGui import QFont, QCursor, QColor, QPolygonF

# PowerPoint Imports
from pptx import Presentation
from pptx.util import Inches, Pt

import pyqtgraph as pg
import pyqtgraph.exporters

# ==========================================
# PYQTGRAPH GLOBAL CONFIGURATION (OPENGL ENABLED)
# ==========================================
try:
    import OpenGL
    pg.setConfigOptions(useOpenGL=True)
except ImportError:
    print("PyOpenGL not found. Falling back to CPU.")
    pg.setConfigOptions(useOpenGL=False)

pg.setConfigOptions(antialias=False) 
pg.setConfigOption('background', '#FFFFFF') 
pg.setConfigOption('foreground', '#5F6368') 

# ==========================================
# CUSTOM GRAPHICS ITEMS
# ==========================================
class TimeAxisItem(pg.AxisItem):
    """Custom X-axis to convert relative seconds + Start Timestamp into HH:MM:SS"""
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.start_timestamp = 0.0

    def tickStrings(self, values, scale, spacing):
        strings = []
        for v in values:
            try:
                real_time = self.start_timestamp + v
                dt = datetime.fromtimestamp(real_time)
                strings.append(dt.strftime("%d-%b %H:%M:%S"))
            except (ValueError, OSError, OverflowError, TypeError):
                strings.append("")
        return strings

class SpeechBubbleAnnotation(pg.GraphicsObject):
    """Liquid Glass Hover Marker with Selection Highlighting & Editing"""
    def __init__(self, x, y, title, text, start_timestamp, delete_cb, edit_cb, select_cb):
        super().__init__()
        self.x = x
        self.y = y
        self.title = title if title else "Note"
        self.text = text
        self.start_timestamp = start_timestamp
        
        self.delete_callback = delete_cb
        self.edit_callback = edit_cb
        self.select_callback = select_cb
        
        self.is_hovered = False
        self.is_selected = False 
        
        self.setPos(x, y)
        self.setFlag(self.ItemIgnoresTransformations) 
        self.setAcceptHoverEvents(True)
        self.setAcceptedMouseButtons(Qt.LeftButton | Qt.RightButton) 
        self.setZValue(1000) 
        
        self.proxy = QGraphicsProxyWidget(self)
        self.text_browser = QTextBrowser()
        self.text_browser.setOpenExternalLinks(True)
        self.text_browser.setReadOnly(True)
        
        self.update_html_content()
        
        # UI Colors matched to new Emerald Theme
        self.text_browser.setStyleSheet("""
            QTextBrowser {
                background-color: rgba(255, 255, 255, 245);
                border: 1px solid rgba(180, 180, 180, 200);
                border-radius: 8px;
            }
            QScrollBar:vertical { width: 8px; background: transparent; }
            QScrollBar::handle:vertical { background: #bcc6cc; border-radius: 4px; }
        """)
        
        self.proxy.setWidget(self.text_browser)
        self.proxy.setPos(-120, -180) 
        self.text_browser.setFixedSize(250, 150)
        self.proxy.setVisible(False)

    def update_content(self, new_title, new_text):
        self.title = new_title if new_title else "Note"
        self.text = new_text
        self.update_html_content()
        self.update()

    def update_html_content(self):
        try:
            dt_str = datetime.fromtimestamp(self.start_timestamp + self.x).strftime("%d-%m-%Y %H:%M:%S")
        except:
            dt_str = "..."
            
        formatted_text = self.text.replace('\n', '<br>')
        html_content = f"""
        <div style="font-family: 'Segoe UI'; padding: 2px;">
            <div style="color: #059669; font-size: 14px; font-weight: bold;">{self.title}</div>
            <div style="color: #555555; font-size: 11px; margin-top: 5px; margin-bottom: 8px;">
                <b>Time:</b> {dt_str} &nbsp;|&nbsp; <b>Value:</b> {self.y:.3f}
            </div>
            <div style="color: #202124; font-size: 12px; line-height: 1.4;">
                {formatted_text}
            </div>
        </div>
        """
        self.text_browser.setHtml(html_content)

    def boundingRect(self):
        if self.is_hovered: return QRectF(-130, -190, 270, 210)
        return QRectF(-15, -20, 180, 40)

    def paint(self, p, option, widget):
        p.setRenderHint(p.Antialiasing) 
        
        if self.is_selected:
            p.setBrush(QColor(5, 150, 105, 20)) # Soft Emerald
            p.setPen(pg.mkPen(color='#059669', width=1.5, style=Qt.DashLine))
            p.drawRoundedRect(QRectF(-12, -20, 175, 40), 6, 6)
        
        p.setBrush(QColor('#202124'))
        p.setPen(pg.mkPen('#FFFFFF', width=2))
        p.drawEllipse(QRectF(-6, -6, 12, 12))
        
        p.setPen(QColor('#059669'))
        p.setFont(QFont("Segoe UI", 10, QFont.Bold))
        p.drawText(QRectF(14, -12, 150, 24), Qt.AlignLeft | Qt.AlignVCenter, self.title)
        
        if self.is_hovered:
            p.setBrush(QColor(255, 255, 255, 245))
            p.setPen(pg.mkPen(QColor(180, 180, 180, 200), width=1))
            triangle = QPolygonF([QPointF(0, -8), QPointF(-10, -30), QPointF(10, -30)])
            p.drawPolygon(triangle)

    def hoverEnterEvent(self, ev):
        self.is_hovered = True
        self.proxy.setVisible(True)
        self.update()

    def hoverLeaveEvent(self, ev):
        self.is_hovered = False
        self.proxy.setVisible(False)
        self.update()

    def mouseClickEvent(self, ev):
        ev.accept()
        self.select_callback(self)
        if ev.button() == Qt.RightButton:
            menu = QMenu()
            menu.setStyleSheet("""
                QMenu { background-color: white; border: 1px solid #DADCE0; border-radius: 6px; padding: 4px; }
                QMenu::item { padding: 8px 24px; font-family: 'Segoe UI'; font-size: 10pt; color: #202124; }
                QMenu::item:selected { background-color: #F1F3F4; border-radius: 4px; }
            """)
            edit_action = menu.addAction("✏️ Edit Note")
            del_action = menu.addAction("🗑️ Delete Note")
            
            try: click_pos = ev.screenPos().toPoint()
            except AttributeError: click_pos = ev.screenPos()
                
            action = menu.exec_(click_pos)
            if action == edit_action: self.edit_callback(self)
            elif action == del_action: self.delete_callback(self)


# ==========================================
# BEAUTIFUL TEXT EDITOR DIALOG (Emerald Theme)
# ==========================================
class AnnotationDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Add / Edit Note")
        self.setMinimumSize(450, 350) 
        self.setStyleSheet("""
            QDialog { background-color: #F8F9FA; }
            QLabel { font-family: 'Segoe UI'; font-weight: bold; color: #059669; font-size: 11pt; margin-top: 4px; }
            QLineEdit, QTextEdit { 
                border: 1px solid #DADCE0; border-radius: 6px; padding: 10px;
                background-color: #FFFFFF; font-family: 'Segoe UI'; font-size: 11pt; color: #202124;
            }
            QLineEdit:focus, QTextEdit:focus { border: 2px solid #059669; padding: 9px; }
            QPushButton { 
                background-color: #059669; color: white; padding: 10px 20px; 
                border-radius: 6px; font-weight: bold; font-family: 'Segoe UI'; font-size: 10pt;
            }
            QPushButton:hover { background-color: #047857; }
            QPushButton#cancelBtn { background-color: #FFFFFF; color: #5F6368; border: 1px solid #DADCE0; }
            QPushButton#cancelBtn:hover { background-color: #F1F3F4; color: #202124; }
        """)
        
        layout = QVBoxLayout(self)
        layout.setContentsMargins(20, 20, 20, 20)
        
        layout.addWidget(QLabel("Title:"))
        self.title_input = QLineEdit()
        self.title_input.setPlaceholderText("e.g. Peak Anomaly")
        layout.addWidget(self.title_input)
        
        layout.addWidget(QLabel("Detailed Notes:"))
        self.text_input = QTextEdit()
        self.text_input.setPlaceholderText("Enter detailed observations or notes here...\nPress Enter to drop to a new line.")
        layout.addWidget(self.text_input)
        
        btn_layout = QHBoxLayout()
        save_btn = QPushButton("Save Note")
        save_btn.clicked.connect(self.accept)
        cancel_btn = QPushButton("Cancel")
        cancel_btn.setObjectName("cancelBtn")
        cancel_btn.clicked.connect(self.reject)
        
        btn_layout.addStretch()
        btn_layout.addWidget(cancel_btn)
        btn_layout.addWidget(save_btn)
        layout.addLayout(btn_layout)

# ==========================================
# FILE PARSING ENGINE (HYBRID TIME ENGINE + DUAL AXIS EXTRACTION)
# ==========================================
class FileMonitorThread(QThread):
    newData = pyqtSignal(np.ndarray, np.ndarray, np.ndarray)
    
    def __init__(self):
        super().__init__()
        self.file_path = None
        self.running = False
        self.current_index = 0
        self.calculated_sample_time = 1.13
        self.start_timestamp = time.time()
        self.ingestion_delay_ms = 100 

    def set_delay(self, ms_val):
        self.ingestion_delay_ms = ms_val

    def start_reading(self, file_path):
        self.stop()
        self.file_path = file_path
        self.current_index = 0
        self.calculated_sample_time = self._calculate_sample_time(file_path)
        self.running = True
        self.start()

    def _calculate_sample_time(self, filepath):
        start_time, end_time, valid_rows, is_live = None, None, 0, True
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                for line in f:
                    line = line.strip()
                    if not line: continue
                    match = re.search(r"(\d{2}-\d{2}-\d{4})\s*-\s*(\d{2}:\d{2}:\d{2})", line)
                    if match:
                        dt = datetime.strptime(f"{match.group(1)} {match.group(2)}", "%d-%m-%Y %H:%M:%S")
                        if start_time is None:
                            start_time = dt
                            self.start_timestamp = start_time.timestamp()
                        else: end_time = dt
                    elif "End log file" in line: is_live = False
                    else:
                        cols = line.split()
                        if len(cols) >= 2:
                            try:
                                float(cols[1].replace('\x00', '').replace('\ufeff', ''))
                                valid_rows += 1
                            except ValueError: pass

            if start_time and valid_rows > 0:
                if not is_live and end_time: 
                    return abs((end_time - start_time).total_seconds()) / valid_rows
                else:
                    calc = abs((datetime.now() - start_time).total_seconds()) / valid_rows
                    return calc if calc > 0 else 1.13
        except Exception as e: print("Calculation error:", e)
        return 1.13 

    def stop(self):
        self.running = False
        self.quit()
        self.wait()

    def run(self):
        if not self.file_path: return
        try:
            with open(self.file_path, "r", encoding='utf-8', errors='ignore') as f:
                x_batch, y_batch, t_batch = [], [], []
                partial_line_buffer = ""
                is_live_tracking = False  # HYBRID ENGINE: Tracks if we hit the live data stream
                
                while self.running:
                    chunk = f.read(65536) 
                    
                    if not chunk:
                        is_live_tracking = True # HYBRID ENGINE: Reached the bottom, switch to PC Clock
                        if x_batch:
                            self.newData.emit(np.array(x_batch, dtype=float), np.array(y_batch, dtype=float), np.array(t_batch, dtype=float))
                            x_batch, y_batch, t_batch = [], [], []
                        time.sleep(max(0.01, self.ingestion_delay_ms / 1000.0))
                        continue

                    lines = (partial_line_buffer + chunk).split('\n')
                    partial_line_buffer = lines.pop() 
                    
                    for line in lines:
                        line_str = line.strip()
                        if not line_str: continue
                        
                        cols = line_str.split()
                        
                        if len(cols) < 2:
                            continue
                            
                        try:
                            # Primary Y1 Extraction (Index 1) perfectly parsing + and - values
                            clean_val = cols[1].replace('\x00', '').replace('\ufeff', '')
                            y_val = float(clean_val)
                            
                            # --- THE HYBRID NEVER-FAIL TIME ENGINE ---
                            if is_live_tracking:
                                # Lock to exact PC atomic clock for live fluctuating 600-700ms data
                                x_val = time.time() - self.start_timestamp
                            else:
                                # Use mathematical spacing to instantly load days of historical backlog
                                x_val = self.current_index * self.calculated_sample_time
                            # -----------------------------------------
                            
                            # Y2 Extraction (H2 Conc %, last column index)
                            t_val = np.nan
                            try:
                                h2_str = cols[-1].replace('\x00', '').replace('\ufeff', '').replace('%', '')
                                t_val = float(h2_str)
                            except ValueError:
                                pass

                            x_batch.append(x_val)
                            y_batch.append(y_val)
                            t_batch.append(t_val)
                            self.current_index += 1
                        except ValueError:
                            pass
                            
                    if len(x_batch) >= 1000:
                        self.newData.emit(np.array(x_batch, dtype=float), np.array(y_batch, dtype=float), np.array(t_batch, dtype=float))
                        x_batch, y_batch, t_batch = [], [], []

        except Exception as e: print("File reading error:", e)


# ==========================================
# MAIN APPLICATION GUI (Emerald/Orange Theme)
# ==========================================
class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("MNST Smart Probe IPU")
        self.resize(1450, 850)
        self.setMinimumSize(1100, 650)
        
        self.filename_str = "No File"
        
        self.capacity = 1_000_000 
        self.x_data = np.zeros(self.capacity, dtype=float)
        self.y_data = np.zeros(self.capacity, dtype=float)
        self.t_data = np.zeros(self.capacity, dtype=float)
        self.data_len = 0
        
        self.paused = False
        self.notes_visible = True
        self.y2_enabled = True
        self.annotations = [] 

        self.setup_ui()
        
        self.stats_timer = QTimer()
        self.stats_timer.timeout.connect(self.update_visible_statistics)
        self.stats_timer.start(250) 

        self.monitor = FileMonitorThread()
        self.monitor.newData.connect(self.update_plot_data)

    def _create_shadow(self):
        shadow = QGraphicsDropShadowEffect()
        shadow.setBlurRadius(15)
        shadow.setColor(QColor(0, 0, 0, 20))
        shadow.setOffset(0, 4)
        return shadow

    def setup_ui(self):
        app_font = QFont("Segoe UI", 10)
        QApplication.setFont(app_font)

        central_widget = QWidget()
        central_widget.setStyleSheet("background-color: #F8F9FA;") 
        self.setCentralWidget(central_widget)
        
        main_layout = QHBoxLayout(central_widget)
        main_layout.setContentsMargins(16, 16, 16, 16)
        main_layout.setSpacing(10)

        # ------------------------------------------
        # LEFT PANEL (Controls & Stats)
        # ------------------------------------------
        self.left_card = QFrame()
        self.left_card.setMinimumWidth(0)
        self.left_card.setMaximumWidth(340)
        self.left_card.setSizePolicy(QSizePolicy.Fixed, QSizePolicy.Expanding)
        self.left_card.setStyleSheet("""
            QFrame { background-color: #FFFFFF; border-radius: 12px; border: 1px solid #E8EAED; }
            QLabel { border: none; color: #202124; }
        """)
        self.left_card.setGraphicsEffect(self._create_shadow()) 
        
        card_layout = QVBoxLayout(self.left_card)
        card_layout.setContentsMargins(20, 20, 20, 20)
        
        # Reduced spacing to fit the new fixed title
        card_layout.setSpacing(12)

        # --- FIXED ANALYZER TITLE ---
        self.main_session_title = QLineEdit("Smart Probe")
        self.main_session_title.setReadOnly(True)
        self.main_session_title.setFont(QFont("Segoe UI", 14, QFont.Bold))
        self.main_session_title.setStyleSheet("""
            background: transparent;
            border: none;
            color: #059669;
        """)
        card_layout.addWidget(self.main_session_title)

        divider_top = QFrame()
        divider_top.setFrameShape(QFrame.HLine)
        divider_top.setStyleSheet("background-color: #E8EAED; border: none;")
        divider_top.setFixedHeight(1)
        card_layout.addWidget(divider_top)
        # ----------------------------

        lbl_stats_title = QLabel("Visible Statistics")
        lbl_stats_title.setFont(QFont("Segoe UI", 10, QFont.Bold))
        lbl_stats_title.setStyleSheet("color: #059669;") # Emerald
        card_layout.addWidget(lbl_stats_title)

        stats_form = QFormLayout()
        stats_form.setVerticalSpacing(10)
        
        self.lbl_points = QLabel("0")
        self.lbl_max = QLabel("0.000")
        self.lbl_min = QLabel("0.000")
        self.lbl_mean = QLabel("0.000")
        self.lbl_std = QLabel("0.000")

        font_metric = QFont("Segoe UI", 10, QFont.Bold)
        for val_lbl in [self.lbl_points, self.lbl_max, self.lbl_min, self.lbl_mean, self.lbl_std]:
            val_lbl.setFont(font_metric)

        stats_form.addRow("Visible Points:", self.lbl_points)
        stats_form.addRow("Max Value:", self.lbl_max)
        stats_form.addRow("Min Value:", self.lbl_min)
        stats_form.addRow("Mean:", self.lbl_mean)
        stats_form.addRow("Std Dev:", self.lbl_std)
        card_layout.addLayout(stats_form)

        divider = QFrame()
        divider.setFrameShape(QFrame.HLine)
        divider.setStyleSheet("background-color: #E8EAED; border: none;")
        divider.setFixedHeight(1)
        card_layout.addWidget(divider)

        # --- GRAPH LABELS & Y2 ENABLE EDITOR ---
        lbl_labels_title = QLabel("Graph Settings")
        lbl_labels_title.setFont(QFont("Segoe UI", 10, QFont.Bold))
        lbl_labels_title.setStyleSheet("color: #059669;") # Emerald
        card_layout.addWidget(lbl_labels_title)

        labels_form = QFormLayout()
        labels_form.setVerticalSpacing(8)

        input_style = """
            QDateTimeEdit, QDoubleSpinBox, QLineEdit {
                background-color: #F1F3F4; border-radius: 6px; padding: 4px; border: 1px solid transparent;
                color: #202124; font-family: 'Segoe UI'; font-size: 9pt;
            }
            QDateTimeEdit:focus, QDoubleSpinBox:focus, QLineEdit:focus { border: 1px solid #059669; background-color: #FFFFFF; }
        """

        self.title_edit = QLineEdit("Real-Time Multiday Sensor Plot")
        self.xaxis_edit = QLineEdit("Time (Real-World)")
        self.yaxis_edit = QLineEdit("EMF Value (mV)")
        
        for edit in [self.title_edit, self.xaxis_edit, self.yaxis_edit]:
            edit.setStyleSheet(input_style)
            edit.textChanged.connect(self.update_graph_labels)

        labels_form.addRow("Title:", self.title_edit)
        labels_form.addRow("X Axis:", self.xaxis_edit)
        labels_form.addRow("Y Axis:", self.yaxis_edit)
        
        self.chk_y2 = QCheckBox("Enable Y2 Axis (H2 Conc %)")
        self.chk_y2.setChecked(True)
        self.chk_y2.stateChanged.connect(self.toggle_y2)
        labels_form.addRow(self.chk_y2)
        
        card_layout.addLayout(labels_form)

        divider_lbl = QFrame()
        divider_lbl.setFrameShape(QFrame.HLine)
        divider_lbl.setStyleSheet("background-color: #E8EAED; border: none;")
        divider_lbl.setFixedHeight(1)
        card_layout.addWidget(divider_lbl)

        # --- STREAM DELAY ---
        lbl_delay_title = QLabel("Plotting Stream Delay")
        lbl_delay_title.setFont(QFont("Segoe UI", 10, QFont.Bold))
        lbl_delay_title.setStyleSheet("color: #059669;") # Emerald
        card_layout.addWidget(lbl_delay_title)

        self.lbl_delay_val = QLabel("Buffer Delay: 150 ms (Safe)")
        self.lbl_delay_val.setStyleSheet("color: #5F6368; font-size: 9pt;")
        card_layout.addWidget(self.lbl_delay_val)

        self.delay_slider = QSlider(Qt.Horizontal)
        self.delay_slider.setRange(0, 1000)
        self.delay_slider.setValue(150)
        self.delay_slider.setTickInterval(100)
        self.delay_slider.setTickPosition(QSlider.TicksBelow)
        self.delay_slider.valueChanged.connect(self.on_delay_changed)
        card_layout.addWidget(self.delay_slider)

        divider2 = QFrame()
        divider2.setFrameShape(QFrame.HLine)
        divider2.setStyleSheet("background-color: #E8EAED; border: none;")
        divider2.setFixedHeight(1)
        card_layout.addWidget(divider2)

        lbl_scale_title = QLabel("Axis Scaling")
        lbl_scale_title.setFont(QFont("Segoe UI", 10, QFont.Bold))
        lbl_scale_title.setStyleSheet("color: #059669;") # Emerald
        card_layout.addWidget(lbl_scale_title)

        scale_form = QFormLayout()
        scale_form.setVerticalSpacing(8)
        
        self.x_min_edit = QDateTimeEdit()
        self.x_max_edit = QDateTimeEdit()
        for edit in [self.x_min_edit, self.x_max_edit]:
            edit.setDisplayFormat("dd-MM-yyyy HH:mm:ss")
            edit.setStyleSheet(input_style)
        
        self.y_min_edit = QDoubleSpinBox()
        self.y_max_edit = QDoubleSpinBox()
        for edit in [self.y_min_edit, self.y_max_edit]:
            edit.setRange(-99999, 99999)
            edit.setDecimals(3)
            edit.setStyleSheet(input_style)

        scale_form.addRow("X Min:", self.x_min_edit)
        scale_form.addRow("X Max:", self.x_max_edit)
        scale_form.addRow("Y Min (mV):", self.y_min_edit)
        scale_form.addRow("Y Max (mV):", self.y_max_edit)
        card_layout.addLayout(scale_form)

        btn_apply_scale = QPushButton("Apply Manual Limits")
        btn_apply_scale.setStyleSheet("QPushButton { background-color: #F1F3F4; color: #059669; padding: 8px; border-radius: 6px; font-weight: bold; } QPushButton:hover { background-color: #D1FAE5; }")
        btn_apply_scale.setCursor(QCursor(Qt.PointingHandCursor))
        btn_apply_scale.clicked.connect(self.apply_manual_scale)
        card_layout.addWidget(btn_apply_scale)

        btn_auto_scale = QPushButton("Auto-Scale Graph")
        btn_auto_scale.setStyleSheet("QPushButton { background-color: #059669; color: white; padding: 8px; border-radius: 6px; font-weight: bold; } QPushButton:hover { background-color: #047857; }")
        btn_auto_scale.setCursor(QCursor(Qt.PointingHandCursor))
        btn_auto_scale.clicked.connect(self.enable_auto_scale)
        card_layout.addWidget(btn_auto_scale)

        card_layout.addStretch()
        main_layout.addWidget(self.left_card, stretch=0)

        # ------------------------------------------
        # RIGHT PANEL 
        # ------------------------------------------
        right_panel = QWidget()
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(0, 0, 0, 0)
        right_layout.setSpacing(16)
        
        top_bar_container = QWidget()
        top_bar = QHBoxLayout(top_bar_container)
        top_bar.setContentsMargins(0, 0, 0, 0)
        top_bar.setSpacing(10)
        
        btn_style = "QPushButton { background-color: #FFFFFF; color: #5F6368; padding: 8px 14px; border-radius: 6px; border: 1px solid #DADCE0; font-weight: bold; } QPushButton:hover { background-color: #F8F9FA; color: #202124; }"

        # --- PANEL TOGGLE BUTTON ---
        self.btn_toggle_panel = QPushButton("◀")
        self.btn_toggle_panel.setFixedWidth(40)
        self.btn_toggle_panel.setStyleSheet(btn_style)
        self.btn_toggle_panel.setCursor(QCursor(Qt.PointingHandCursor))
        self.btn_toggle_panel.clicked.connect(self.toggle_left_panel)
        top_bar.addWidget(self.btn_toggle_panel)

        self.file_label = QLabel("File: None")
        self.file_label.setFont(QFont("Segoe UI", 12, QFont.Bold))
        self.file_label.setStyleSheet("color: #202124; margin-left: 10px;")
        top_bar.addWidget(self.file_label)
        
        # Emerald highlighted readout
        self.live_emf_display = QLabel("0.000 mV")
        self.live_emf_display.setFont(QFont("Segoe UI", 14, QFont.Bold))
        self.live_emf_display.setStyleSheet("background-color: #D1FAE5; color: #059669; padding: 6px 14px; border-radius: 8px; border: 1px solid #A7F3D0;")
        top_bar.addWidget(self.live_emf_display)
        
        top_bar.addStretch()

        # --- MOUSE MODE TOGGLE (ZOOM vs PAN) ---
        self.btn_mouse_mode = QPushButton("✋ Pan Mode")
        self.btn_mouse_mode.setStyleSheet(btn_style)
        self.btn_mouse_mode.setCursor(QCursor(Qt.PointingHandCursor))
        self.btn_mouse_mode.clicked.connect(self.toggle_mouse_mode)
        top_bar.addWidget(self.btn_mouse_mode)

        self.btn_toggle_notes = QPushButton("👁️ Hide Notes")
        self.btn_toggle_notes.setStyleSheet(btn_style)
        self.btn_toggle_notes.setCursor(QCursor(Qt.PointingHandCursor))
        self.btn_toggle_notes.clicked.connect(self.toggle_notes_visibility)
        top_bar.addWidget(self.btn_toggle_notes)

        self.btn_load = QPushButton("Load File")
        self.btn_load.setStyleSheet("QPushButton { background-color: #059669; color: white; padding: 8px 14px; border-radius: 6px; font-weight: bold; } QPushButton:hover { background-color: #047857; }")
        self.btn_load.setCursor(QCursor(Qt.PointingHandCursor))
        self.btn_load.clicked.connect(self.load_file)
        top_bar.addWidget(self.btn_load)

        self.btn_pause = QPushButton("Pause")
        self.btn_pause.setStyleSheet("QPushButton { background-color: #EA580C; color: white; padding: 8px 14px; border-radius: 6px; font-weight: bold; } QPushButton:hover { background-color: #C2410C; }")
        self.btn_pause.setCursor(QCursor(Qt.PointingHandCursor))
        self.btn_pause.clicked.connect(self.toggle_pause)
        top_bar.addWidget(self.btn_pause)

        self.btn_refresh = QPushButton("Refresh")
        self.btn_refresh.setStyleSheet(btn_style)
        self.btn_refresh.setCursor(QCursor(Qt.PointingHandCursor))
        self.btn_refresh.clicked.connect(self.refresh_plot)
        top_bar.addWidget(self.btn_refresh)

        self.btn_save_png = QPushButton("Export PNG")
        self.btn_save_png.setStyleSheet(btn_style)
        self.btn_save_png.setCursor(QCursor(Qt.PointingHandCursor))
        self.btn_save_png.clicked.connect(self.save_png)
        top_bar.addWidget(self.btn_save_png)
        
        self.btn_package = QPushButton("Export Package")
        self.btn_package.setStyleSheet("QPushButton { background-color: #F59E0B; color: white; padding: 8px 14px; border-radius: 6px; font-weight: bold; } QPushButton:hover { background-color: #D97706; }")
        self.btn_package.setCursor(QCursor(Qt.PointingHandCursor))
        self.btn_package.clicked.connect(self.export_analysis_package)
        top_bar.addWidget(self.btn_package)

        right_layout.addWidget(top_bar_container, stretch=0)

        # --- Graph Card ---
        graph_card = QFrame()
        graph_card.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        graph_card.setStyleSheet("background-color: #FFFFFF; border-radius: 12px; border: 2px solid #DADCE0;")
        
        graph_layout = QVBoxLayout(graph_card)
        graph_layout.setContentsMargins(12, 12, 12, 12)

        self.time_axis = TimeAxisItem(orientation='bottom')
        self.graphWidget = pg.PlotWidget(axisItems={'bottom': self.time_axis})
        self.graphWidget.setStyleSheet("border: none; background: #FFFFFF;") 
        
        self.graphWidget.setMenuEnabled(False) 
        self.graphWidget.getAxis('bottom').enableAutoSIPrefix(False)
        self.graphWidget.setMouseEnabled(x=True, y=True)
        
        self.graphWidget.showGrid(x=True, y=True, alpha=0.3)
        self.graphWidget.getAxis('bottom').setPen(pg.mkPen(color='#DADCE0', width=1))
        self.graphWidget.getAxis('left').setPen(pg.mkPen(color='#DADCE0', width=1))
        
        # Bold Black Axis Text
        tick_font = QFont("Segoe UI", 9, QFont.Bold)
        self.graphWidget.getAxis('bottom').setTickFont(tick_font)
        self.graphWidget.getAxis('bottom').setTextPen('#000000')
        
        self.graphWidget.getAxis('left').setTickFont(tick_font)
        self.graphWidget.getAxis('left').setTextPen('#000000')
        
        # Dual Axis Setup (Y2 for H2 Conc)
        self.p2 = pg.ViewBox()
        self.p2.setMenuEnabled(False)
        self.graphWidget.scene().addItem(self.p2)
        self.graphWidget.getAxis('right').linkToView(self.p2)
        self.p2.setXLink(self.graphWidget)
        self.graphWidget.getAxis('right').setLabel('H2 Conc (%)', color='#EA580C', **{'font-size': '11pt', 'bold': True})
        self.graphWidget.getAxis('right').setTickFont(tick_font)
        self.graphWidget.getAxis('right').setTextPen('#000000')

        # Handles resize dynamics accurately mapping Y2 over Y1
        def updateViews():
            self.p2.setGeometry(self.graphWidget.plotItem.vb.sceneBoundingRect())
            self.p2.linkedViewChanged(self.graphWidget.plotItem.vb, self.p2.XAxis)
            
        updateViews()
        self.graphWidget.plotItem.vb.sigResized.connect(updateViews)
        
        self.graphWidget.scene().sigMouseClicked.connect(self.on_graph_clicked)

        # PRIMARY Y1 PLOT (EMF) - EMERALD GREEN
        self.plot_curve = self.graphWidget.plot(
            pen=pg.mkPen('#059669', width=2),
            symbol='o', symbolSize=6, symbolBrush='#059669', symbolPen=None, 
            autoDownsample=True, downsampleMethod='subsample', clipToView=True
        )
        
        # SECONDARY Y2 PLOT (H2 Conc) - DEEP ORANGE
        self.plot_curve_temp = pg.PlotCurveItem(
            pen=pg.mkPen('#EA580C', width=2, style=Qt.DashLine),
            symbol='t', symbolSize=6, symbolBrush='#EA580C', symbolPen=None,
            autoDownsample=True, clipToView=True
        )
        self.p2.addItem(self.plot_curve_temp)
        
        # Crosshairs - Purple for high visibility against green/orange
        self.vLine = pg.InfiniteLine(angle=90, movable=False, pen=pg.mkPen('#7C3AED', width=1.5, style=Qt.DashLine))
        self.hLine = pg.InfiniteLine(angle=0, movable=False, pen=pg.mkPen('#7C3AED', width=1.5, style=Qt.DashLine))
        self.graphWidget.addItem(self.vLine, ignoreBounds=True)
        self.graphWidget.addItem(self.hLine, ignoreBounds=True)
        
        self.tooltip_label = pg.TextItem(text="", color="#202124", fill=pg.mkBrush(255, 255, 255, 220))
        self.graphWidget.addItem(self.tooltip_label, ignoreBounds=True)
        
        self.proxy = pg.SignalProxy(self.graphWidget.scene().sigMouseMoved, rateLimit=30, slot=self.mouse_moved)
        
        graph_layout.addWidget(self.graphWidget)
        right_layout.addWidget(graph_card, stretch=1)
        main_layout.addWidget(right_panel, stretch=1)
        
        self.update_graph_labels()
        self.toggle_y2(Qt.Checked)
        self.setFocusPolicy(Qt.StrongFocus)

    # ==========================================
    # UI CONTROLS & ANIMATIONS
    # ==========================================
    def toggle_left_panel(self):
        """Smoothly animates the left control panel open or closed."""
        self.anim = QPropertyAnimation(self.left_card, b"maximumWidth")
        self.anim.setDuration(250)
        self.anim.setEasingCurve(pg.QtCore.QEasingCurve.InOutQuad)
        
        if self.left_card.width() > 10:
            self.anim.setStartValue(340)
            self.anim.setEndValue(0)
            self.btn_toggle_panel.setText("▶")
        else:
            self.anim.setStartValue(0)
            self.anim.setEndValue(340)
            self.btn_toggle_panel.setText("◀")
        self.anim.start()

    def update_graph_labels(self):
        """Dynamically updates graph labels with bold text from the text boxes"""
        self.graphWidget.setTitle(self.title_edit.text(), color="#202124", size="16pt", bold=True)
        self.graphWidget.setLabel('bottom', self.xaxis_edit.text(), **{'color': '#5F6368', 'font-size': '11pt', 'bold': True})
        self.graphWidget.setLabel('left', self.yaxis_edit.text(), **{'color': '#059669', 'font-size': '11pt', 'bold': True})

    def toggle_y2(self, state):
        self.y2_enabled = state == Qt.Checked
        self.graphWidget.getAxis('right').setVisible(self.y2_enabled)
        self.plot_curve_temp.setVisible(self.y2_enabled)
        self.graphWidget.update()

    def toggle_mouse_mode(self):
        """Toggles between standard panning and box-zooming (screenshot style)"""
        view_box = self.graphWidget.getViewBox()
        if view_box.state['mouseMode'] == pg.ViewBox.RectMode:
            view_box.setMouseMode(pg.ViewBox.PanMode)
            self.btn_mouse_mode.setText("✋ Pan Mode")
            self.btn_mouse_mode.setStyleSheet("QPushButton { background-color: #FFFFFF; color: #5F6368; padding: 8px 14px; border-radius: 6px; border: 1px solid #DADCE0; font-weight: bold; } QPushButton:hover { background-color: #F8F9FA; color: #202124; }")
        else:
            view_box.setMouseMode(pg.ViewBox.RectMode)
            self.btn_mouse_mode.setText("🔍 Box Zoom Mode")
            self.btn_mouse_mode.setStyleSheet("QPushButton { background-color: #D1FAE5; color: #059669; padding: 8px 14px; border-radius: 6px; border: 1px solid #A7F3D0; font-weight: bold; }")

    def on_delay_changed(self, value):
        self.lbl_delay_val.setText(f"Buffer Delay: {value} ms")
        if self.monitor:
            self.monitor.set_delay(value)

    # ==========================================
    # USER INTERACTIONS & ANNOTATIONS
    # ==========================================
    def on_graph_clicked(self, event):
        items_under_mouse = self.graphWidget.scene().items(event.scenePos())
        for item in items_under_mouse:
            if isinstance(item, SpeechBubbleAnnotation) or isinstance(item, QGraphicsProxyWidget):
                return

        if event.double():
            self.enable_auto_scale()
            
        elif event.button() == Qt.LeftButton:
            for ann in self.annotations:
                ann.is_selected = False
                ann.update()
                
        elif event.button() == Qt.RightButton:
            for ann in self.annotations:
                ann.is_selected = False
                ann.update()
                
            pos = event.scenePos()
            if self.graphWidget.sceneBoundingRect().contains(pos):
                mousePoint = self.graphWidget.plotItem.vb.mapSceneToView(pos)
                menu = QMenu()
                menu.setStyleSheet("""
                    QMenu { background-color: white; border: 1px solid #DADCE0; border-radius: 6px; padding: 4px; }
                    QMenu::item { padding: 8px 24px; font-family: 'Segoe UI'; font-size: 10pt; color: #202124; }
                    QMenu::item:selected { background-color: #F1F3F4; border-radius: 4px; }
                """)
                add_action = menu.addAction("📝 Place a Note Here")
                
                try: click_pos = event.screenPos().toPoint()
                except AttributeError: click_pos = event.screenPos()
                
                action = menu.exec_(click_pos)
                if action == add_action:
                    self.create_annotation(mousePoint.x(), mousePoint.y())
                    
    def create_annotation(self, x, y):
        dlg = AnnotationDialog(self)
        if dlg.exec_():
            title = dlg.title_input.text().strip()
            text = dlg.text_input.toPlainText().strip()
            if title or text:
                ann = SpeechBubbleAnnotation(
                    x, y, title, text, 
                    self.monitor.start_timestamp, 
                    self.delete_annotation, 
                    self.edit_annotation, 
                    self.select_annotation
                )
                self.annotations.append(ann)
                self.graphWidget.addItem(ann)
                self.select_annotation(ann)
                if not self.notes_visible:
                    self.toggle_notes_visibility()

    def select_annotation(self, selected_ann):
        for ann in self.annotations:
            ann.is_selected = (ann == selected_ann)
            ann.update()

    def edit_annotation(self, ann):
        dlg = AnnotationDialog(self)
        dlg.title_input.setText(ann.title)
        dlg.text_input.setText(ann.text)
        if dlg.exec_():
            ann.update_content(dlg.title_input.text().strip(), dlg.text_input.toPlainText().strip())

    def delete_annotation(self, ann):
        if ann in self.annotations:
            self.annotations.remove(ann)
            self.graphWidget.removeItem(ann)

    def toggle_notes_visibility(self):
        self.notes_visible = not self.notes_visible
        self.btn_toggle_notes.setText("👁️ Hide Notes" if self.notes_visible else "👁️‍🗨️ Show Notes")
        for ann in self.annotations: ann.setVisible(self.notes_visible)

    # ==========================================
    # FILE & DATA MANAGEMENT
    # ==========================================
    def load_file(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Open Log File", "", "Log/Text Files (*.log *.txt);;All Files (*.*)")
        if file_path:
            self.filename_str = os.path.basename(file_path)
            self.setWindowTitle(f"MNST Plotter - {self.filename_str}")
            self.file_label.setText(f"File: {self.filename_str}")
            
            self._reset_data()
            self.monitor.start_reading(file_path)

    def refresh_plot(self):
        self._reset_data()
        if self.monitor.file_path:
            self.monitor.start_reading(self.monitor.file_path)
            
    def _reset_data(self):
        self.data_len = 0
        self.capacity = 1_000_000 
        self.x_data = np.zeros(self.capacity, dtype=float)
        self.y_data = np.zeros(self.capacity, dtype=float)
        self.t_data = np.zeros(self.capacity, dtype=float)
        self.plot_curve.setData(self.x_data[:self.data_len], self.y_data[:self.data_len])
        self.plot_curve_temp.setData(self.x_data[:self.data_len], self.t_data[:self.data_len])
        for ann in self.annotations:
            self.graphWidget.removeItem(ann)
        self.annotations = []

    def toggle_pause(self):
        self.paused = not self.paused
        if self.paused:
            self.btn_pause.setText("Resume")
            self.btn_pause.setStyleSheet("QPushButton { background-color: #059669; color: white; padding: 8px 14px; border-radius: 6px; font-weight: bold; } QPushButton:hover { background-color: #047857; }")
        else:
            self.btn_pause.setText("Pause")
            self.btn_pause.setStyleSheet("QPushButton { background-color: #EA580C; color: white; padding: 8px 14px; border-radius: 6px; font-weight: bold; } QPushButton:hover { background-color: #C2410C; }")

    # ==========================================
    # SCALING & VISIBLE STATISTICS
    # ==========================================
    def apply_manual_scale(self):
        self.graphWidget.disableAutoRange()
        
        x_min_abs = self.x_min_edit.dateTime().toSecsSinceEpoch()
        x_max_abs = self.x_max_edit.dateTime().toSecsSinceEpoch()
        x_min_rel = x_min_abs - self.monitor.start_timestamp
        x_max_rel = x_max_abs - self.monitor.start_timestamp
        y_min = self.y_min_edit.value()
        y_max = self.y_max_edit.value()
        
        if x_max_rel > x_min_rel: self.graphWidget.setXRange(x_min_rel, x_max_rel, padding=0)
        if y_max > y_min: self.graphWidget.setYRange(y_min, y_max, padding=0)
        self.update_visible_statistics()

    def enable_auto_scale(self):
        self.graphWidget.enableAutoRange(axis=pg.ViewBox.XYAxes, enable=True)
        self.p2.enableAutoRange(axis=pg.ViewBox.YAxis, enable=True)

    def update_visible_statistics(self):
        if self.data_len == 0: 
            return
            
        view_box = self.graphWidget.getViewBox()
        x_min, x_max = view_box.viewRange()[0]
        
        start_idx = np.searchsorted(self.x_data[:self.data_len], x_min)
        end_idx = np.searchsorted(self.x_data[:self.data_len], x_max)
        
        visible_y = self.y_data[start_idx:end_idx]
        
        if len(visible_y) == 0:
            self.lbl_points.setText("0")
            for lbl in [self.lbl_max, self.lbl_min, self.lbl_mean, self.lbl_std]:
                lbl.setText("N/A")
            return
            
        self.lbl_points.setText(f"{len(visible_y):,}")
        self.lbl_max.setText(f"{np.max(visible_y):.3f}")
        self.lbl_min.setText(f"{np.min(visible_y):.3f}")
        self.lbl_mean.setText(f"{np.mean(visible_y):.3f}")
        self.lbl_std.setText(f"{np.std(visible_y):.4f}")

    def update_plot_data(self, new_x_arr, new_y_arr, new_t_arr):
        if self.paused: return
        self.time_axis.start_timestamp = self.monitor.start_timestamp
        
        n_new = len(new_x_arr)
        if n_new == 0: return
        
        if self.data_len + n_new > self.capacity:
            self.capacity = max(self.capacity * 2, self.data_len + n_new)
            cache_x = np.zeros(self.capacity, dtype=float)
            cache_y = np.zeros(self.capacity, dtype=float)
            cache_t = np.zeros(self.capacity, dtype=float)
            cache_x[:self.data_len] = self.x_data[:self.data_len]
            cache_y[:self.data_len] = self.y_data[:self.data_len]
            cache_t[:self.data_len] = self.t_data[:self.data_len]
            self.x_data = cache_x
            self.y_data = cache_y
            self.t_data = cache_t
            
        self.x_data[self.data_len:self.data_len + n_new] = new_x_arr
        self.y_data[self.data_len:self.data_len + n_new] = new_y_arr
        self.t_data[self.data_len:self.data_len + n_new] = new_t_arr
        
        if self.data_len == 0:
            self.enable_auto_scale()
            
        self.data_len += n_new
        
        self.plot_curve.setData(self.x_data[:self.data_len], self.y_data[:self.data_len])
        self.plot_curve_temp.setData(self.x_data[:self.data_len], self.t_data[:self.data_len])
            
        if n_new > 0:
            self.live_emf_display.setText(f"{self.y_data[self.data_len-1]:.3f} mV")

    # ==========================================
    # EXPORT ENGINES
    # ==========================================
    def export_analysis_package(self):
        if self.data_len == 0:
            QMessageBox.warning(self, "No Data", "There is no data to export.")
            return
            
        path, _ = QFileDialog.getSaveFileName(self, "Export Analysis Package", f"{self.filename_str}_Analysis.zip", "ZIP Files (*.zip)")
        if not path: return
        
        try:
            with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zf:
                temp_img = "temp_analysis_plot.png"
                pg.exporters.ImageExporter(self.graphWidget.plotItem).export(temp_img)
                zf.write(temp_img, "Plot_Screenshot.png")
                os.remove(temp_img)
                
                summary = "MNST Analysis Summary Report\n"
                summary += "="*30 + "\n\n"
                summary += f"Source Log File: {self.filename_str}\n"
                summary += f"Total Data Points: {self.data_len:,}\n"
                
                active_y = self.y_data[:self.data_len]
                if self.data_len > 0:
                    summary += f"Overall Max: {np.max(active_y):.3f} mV\n"
                    summary += f"Overall Min: {np.min(active_y):.3f} mV\n"
                    summary += f"Overall Mean: {np.mean(active_y):.3f} mV\n"
                summary += "\n" + "-"*40 + "\n\n"
                
                summary += "USER ANNOTATIONS & MARKERS:\n\n"
                if self.annotations:
                    for idx, ann in enumerate(self.annotations, 1):
                        try:
                            dt_str = datetime.fromtimestamp(self.monitor.start_timestamp + ann.x).strftime("%d-%m-%Y %H:%M:%S")
                        except:
                            dt_str = "Unknown Time"
                        summary += f"Note #{idx}: {ann.title}\n"
                        summary += f"Timestamp:   {dt_str}\n"
                        summary += f"EMF Value:   {ann.y:.3f} mV\n"
                        summary += f"Description: {ann.text}\n\n"
                else:
                    summary += "(No notes were added during this session.)\n"
                    
                zf.writestr("Analysis_Summary.txt", summary)
                
            QMessageBox.information(self, "Success", "Analysis Package successfully exported!")
        except Exception as e:
            QMessageBox.critical(self, "Export Error", f"Failed to export package: {e}")

    def save_png(self):
        if self.data_len == 0: return
        path, _ = QFileDialog.getSaveFileName(self, "Save PNG", self.filename_str + "_Plot.png", "PNG Files (*.png)")
        if path:
            pg.exporters.ImageExporter(self.graphWidget.plotItem).export(path)
            QMessageBox.information(self, "Success", f"Graph saved to:\n{path}")

    def send_to_ppt(self):
        if self.data_len == 0: return
        ppt_file, _ = QFileDialog.getOpenFileName(self, "Select PowerPoint", "", "PowerPoint (*.pptx)")
        if not ppt_file: return

        temp_img = "temp_sensor_plot.png"
        pg.exporters.ImageExporter(self.graphWidget.plotItem).export(temp_img)

        try:
            prs = Presentation(ppt_file)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            p = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9), Inches(0.5)).text_frame.paragraphs[0]
            p.text = self.filename_str
            p.font.size = Pt(24); p.font.bold = True
            slide.shapes.add_picture(temp_img, Inches(0.5), Inches(0.8), width=Inches(9))
            prs.save(ppt_file)
            if os.path.exists(temp_img): os.remove(temp_img)
            QMessageBox.information(self, "Success", "Plot added to PowerPoint successfully!")
        except Exception as e:
            QMessageBox.critical(self, "PPTX Error", f"Failed to modify PowerPoint:\n{e}")

    def mouse_moved(self, evt):
        pos = evt[0]
        if self.graphWidget.sceneBoundingRect().contains(pos):
            mousePoint = self.graphWidget.plotItem.vb.mapSceneToView(pos)
            self.vLine.setPos(mousePoint.x())
            
            # Locate exact point data via O(log N) search
            idx = np.searchsorted(self.x_data[:self.data_len], mousePoint.x())
            if idx < self.data_len:
                exact_y = self.y_data[idx]
                exact_t = self.t_data[idx]
                self.hLine.setPos(exact_y)
                
                try:
                    real_time = self.monitor.start_timestamp + mousePoint.x()
                    time_str = datetime.fromtimestamp(real_time).strftime("%d-%b %H:%M:%S")
                except Exception: 
                    time_str = "..."
                
                self.tooltip_label.setPos(mousePoint.x(), exact_y)
                if self.y2_enabled and not np.isnan(exact_t):
                    self.tooltip_label.setText(f" {time_str} | Y1: {exact_y:.3f} mV | H2 Conc: {exact_t:.4f} % ")
                else:
                    self.tooltip_label.setText(f" {time_str} | {exact_y:.3f} mV ")

    def keyPressEvent(self, event):
        view_box = self.graphWidget.getViewBox()
        x_range, y_range = view_box.viewRange()
        dx, dy = (x_range[1] - x_range[0]) * 0.05, (y_range[1] - y_range[0]) * 0.05

        if event.key() == Qt.Key_Left: view_box.translateBy(x=-dx, y=0)
        elif event.key() == Qt.Key_Right: view_box.translateBy(x=dx, y=0)
        elif event.key() == Qt.Key_Up: view_box.translateBy(x=0, y=dy)
        elif event.key() == Qt.Key_Down: view_box.translateBy(x=0, y=-dy)
        else: super().keyPressEvent(event)

    def closeEvent(self, event):
        if self.monitor and self.monitor.isRunning():
            self.monitor.stop()
        event.accept()
        QApplication.instance().quit()

if __name__ == "__main__":
    app = QApplication(sys.argv)
    app.setStyle("Fusion") 
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())